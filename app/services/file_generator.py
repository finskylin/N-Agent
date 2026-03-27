
import os
import uuid
import json
import socket
import time
from pathlib import Path
from typing import Dict, List, Any, Optional, Union
from urllib.parse import quote, parse_qs, unquote, urlparse
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import io
import base64
from PIL import Image
import matplotlib.pyplot as plt  # matplotlib需要先安装
import matplotlib.font_manager as fm
import pandas as pd
import numpy as np
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib import colors


# ============================================================
# 动态 IP 获取（用于生成对外可访问的下载链接）
# ============================================================
_cached_local_ip: str = None
_cached_local_ip_time: float = 0
DEFAULT_LOCAL_OBJECT_STORE_DIR = Path("app/data/object_storage")
DEFAULT_GENERATED_FILES_DIR = Path("app/data/generated")


def get_local_ip() -> str:
    """
    动态获取宿主机的局域网 IP 地址

    Docker 容器内优先通过 host.docker.internal 解析宿主机真实 IP，
    确保生成的 URL 可被同局域网设备（如钉钉客户端）访问。
    """
    # 方法 0: 从 HOST_LAN_IP 环境变量获取（docker-compose 启动时自动注入）
    host_lan_ip = os.environ.get("HOST_LAN_IP", "")
    if host_lan_ip and not host_lan_ip.startswith("127."):
        return host_lan_ip

    # 方法 0.5: Docker 容器内通过 host.docker.internal 获取宿主机 IP
    # 注意：Docker Desktop macOS 下返回的是 VM 网关 IP（192.168.65.x），
    # 同局域网设备可能无法访问。优先使用 HOST_LAN_IP。
    try:
        host_ip = socket.gethostbyname("host.docker.internal")
        if host_ip and not host_ip.startswith("127.") and not host_ip.startswith("192.168.65."):
            return host_ip
    except (socket.gaierror, OSError):
        pass

    try:
        # 方法 1: 通过 UDP socket 连接外部地址获取本机出口 IP
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.settimeout(0.1)
        s.connect(("8.8.8.8", 80))
        local_ip = s.getsockname()[0]
        s.close()
        return local_ip
    except Exception:
        pass

    try:
        # 方法 2: 通过 hostname 获取
        import subprocess
        result = subprocess.run(
            ["hostname", "-I"],
            capture_output=True,
            text=True,
            timeout=2
        )
        if result.returncode == 0 and result.stdout.strip():
            ips = result.stdout.strip().split()
            if ips:
                return ips[0]
    except Exception:
        pass

    try:
        # 方法 3: 通过 ip addr 命令获取
        import subprocess
        result = subprocess.run(
            ["ip", "addr", "show"],
            capture_output=True,
            text=True,
            timeout=2
        )
        if result.returncode == 0:
            import re
            matches = re.findall(r'inet (\d+\.\d+\.\d+\.\d+)', result.stdout)
            for ip in matches:
                if not ip.startswith('127.'):
                    return ip
    except Exception:
        pass

    # 方法 4: 返回 localhost 作为最后手段
    return "127.0.0.1"


def get_external_host() -> str:
    """
    获取对外访问的 IP 地址

    优先级：
    1. AGENT_EXTERNAL_HOST 环境变量（最高优先级）
    2. 动态检测本机 IP（缓存 5 分钟）
    """
    global _cached_local_ip, _cached_local_ip_time

    # 1. 优先使用环境变量
    external_host = os.environ.get("AGENT_EXTERNAL_HOST")
    if external_host:
        return external_host

    # 2. 使用缓存的 IP（5 分钟内有效）
    current_time = time.time()
    if _cached_local_ip and (current_time - _cached_local_ip_time) < 300:
        return _cached_local_ip

    # 3. 动态获取本机 IP
    _cached_local_ip = get_local_ip()
    _cached_local_ip_time = current_time
    print(f"[FileGenerator] Detected external host: {_cached_local_ip}")

    return _cached_local_ip


def get_local_object_store_dir() -> Path:
    """获取本地对象存储根目录。"""
    configured = os.environ.get("LOCAL_OBJECT_STORE_DIR", "").strip()
    return Path(configured) if configured else DEFAULT_LOCAL_OBJECT_STORE_DIR


def get_generated_files_dir() -> Path:
    """获取可下载生成文件目录。"""
    configured = os.environ.get("LOCAL_GENERATED_FILES_DIR", "").strip()
    return Path(configured) if configured else DEFAULT_GENERATED_FILES_DIR


def get_public_base_url() -> str:
    """获取下载链接的公开访问前缀。"""
    public_base_url = os.environ.get("AGENT_PUBLIC_BASE_URL", "").strip()
    if public_base_url:
        return public_base_url.rstrip("/")
    external_host = get_external_host()
    port = get_service_port()
    return f"http://{external_host}:{port}"


def get_public_storage_url() -> str:
    """
    获取文件存储对外可访问的 base URL。

    优先级：
    1. MINIO_PUBLIC_URL 环境变量（历史兼容，非空时使用）
    2. 动态获取宿主机真实 IP + 存储端口映射
    """
    storage_public_url = os.environ.get("MINIO_PUBLIC_URL", "")
    if storage_public_url:
        return storage_public_url.rstrip("/")

    host_ip = get_local_ip()
    storage_port = os.environ.get("MINIO_API_PORT", "19000")
    return f"http://{host_ip}:{storage_port}"


def get_service_port() -> int:
    """获取服务端口"""
    # 优先使用环境变量
    port_str = os.environ.get("AGENT_SERVICE_PORT", os.environ.get("PORT", "8000"))
    try:
        return int(port_str)
    except ValueError:
        return 8000


def _normalize_object_name(object_name: str) -> str:
    normalized = str(object_name or "").replace("\\", "/").lstrip("/")
    parts = [part for part in normalized.split("/") if part not in {"", ".", ".."}]
    return "/".join(parts)


def _allowed_download_roots() -> Dict[str, Path]:
    return {
        "object_storage": get_local_object_store_dir().resolve(),
        "generated": get_generated_files_dir().resolve(),
    }


def _resolve_safe_download_token(path: Union[str, Path]) -> str:
    file_path = Path(path)
    resolved_path = file_path.resolve() if file_path.is_absolute() else (Path.cwd() / file_path).resolve()
    for prefix, root in _allowed_download_roots().items():
        try:
            relative = resolved_path.relative_to(root)
            normalized_relative = _normalize_object_name(relative.as_posix())
            if not normalized_relative:
                raise ValueError(f"文件路径不是可下载文件: {resolved_path}")
            return f"{prefix}/{normalized_relative}"
        except ValueError:
            continue
    raise ValueError(f"文件路径不在允许下载的目录内: {resolved_path}")


def resolve_local_path_from_download_value(value: str) -> Optional[Path]:
    """从下载 token、下载 URL 或兼容旧 path 解析到本地文件路径。"""
    if not value:
        return None

    raw_value = str(value).strip()
    if not raw_value:
        return None

    parsed = urlparse(raw_value)
    if parsed.scheme and parsed.path == "/api/files/download":
        path_values = parse_qs(parsed.query).get("path", [])
        if not path_values:
            return None
        raw_value = unquote(path_values[0]).strip()

    normalized = raw_value.replace("\\", "/").strip()
    if not normalized:
        return None

    roots = _allowed_download_roots()
    for prefix, root in roots.items():
        prefix_marker = f"{prefix}/"
        if normalized.startswith(prefix_marker):
            relative = _normalize_object_name(normalized[len(prefix_marker):])
            if not relative:
                return None
            return (root / relative).resolve()

    file_path = Path(normalized)
    if not file_path.is_absolute():
        file_path = (Path.cwd() / file_path).resolve()
    else:
        file_path = file_path.resolve()

    for root in roots.values():
        try:
            file_path.relative_to(root)
            return file_path
        except ValueError:
            continue
    return None


def build_download_url_for_path(path: Union[str, Path]) -> str:
    """为任意本地文件路径生成统一下载 URL。"""
    download_token = _resolve_safe_download_token(path)
    encoded_path = quote(download_token, safe="/:_-.()")
    return f"{get_public_base_url()}/api/files/download?path={encoded_path}"


def save_bytes_to_local_storage(
    file_data: bytes,
    bucket_name: str,
    object_name: str,
) -> Dict[str, str]:
    """
    将二进制内容保存到本地对象存储，目录层级与 bucket/object_name 保持一致。
    """
    safe_bucket = _normalize_object_name(bucket_name) or "uploads"
    safe_object_name = _normalize_object_name(object_name)
    if not safe_object_name:
        raise ValueError("object_name 不能为空")
    relative_path = get_local_object_store_dir() / safe_bucket / safe_object_name
    relative_path.parent.mkdir(parents=True, exist_ok=True)
    relative_path.write_bytes(file_data)
    return {
        "bucket": safe_bucket,
        "object_name": safe_object_name,
        "local_path": str(relative_path),
        "download_url": build_download_url_for_path(relative_path),
    }


def resolve_local_path_from_download_url(url: str) -> Optional[Path]:
    """从 /api/files/download 链接反解到本地文件路径。"""
    return resolve_local_path_from_download_value(url)

# 配置matplotlib中文字体
def configure_matplotlib_chinese():
    """配置matplotlib支持中文字体"""
    try:
        # 尝试设置系统中文字体
        chinese_fonts = [
            'WenQuanYi Zen Hei', # 文泉驿正黑 (Container present)
            'Zen Hei',
            'Liberation Sans',   # Fallback English
            'FreeSans',          # Fallback English
            'Heiti TC',  # 黑体-繁
            'Heiti SC',  # 黑体-简
            'PingFang SC',  # 苹方-简
            'PingFang TC',  # 苹方-繁
            'Microsoft YaHei',  # 微软雅黑
            'SimHei',  # 黑体
            'SimSun',  # 宋体
            'DejaVu Sans'  # 备用
        ]

        for font_name in chinese_fonts:
            try:
                font = fm.FontProperties(family=font_name)
                # 测试字体是否可用
                plt.rcParams['font.sans-serif'] = [font_name]
                plt.rcParams['axes.unicode_minus'] = False
                print(f"[FileGenerator] Successfully configured matplotlib with font: {font_name}")
                return
            except:
                continue

        print("[FileGenerator] Warning: No Chinese font found, using default font")
    except Exception as e:
        print(f"[FileGenerator] Warning: Failed to configure Chinese font: {e}")

# 初始化时配置中文字体
configure_matplotlib_chinese()

class FileGenerator:
    """
    Unified File Generation Service
    Supports: docx, xlsx, pptx, pdf
    Returns: Download URL and file metadata
    """
    
    # Base directory for generated files
    BASE_DIR = get_generated_files_dir()
    
    @classmethod
    def _ensure_dir(cls):
        cls.BASE_DIR = get_generated_files_dir()
        if not cls.BASE_DIR.exists():
            cls.BASE_DIR.mkdir(parents=True, exist_ok=True)

    @classmethod
    def _get_download_url(cls, filename: str) -> str:
        """
        生成完整的对外可访问下载链接

        返回格式: http://{external_host}:{port}/api/files/download?path=data/generated/{filename}

        环境变量配置:
        - AGENT_EXTERNAL_HOST: 指定外部访问地址（可选，默认自动检测）
        - AGENT_SERVICE_PORT: 指定服务端口（可选，默认 8000）
        """
        cls._ensure_dir()
        return build_download_url_for_path(cls.BASE_DIR / filename)

    @classmethod
    def generate_docx(cls, data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Generate Word Document
        Input data structure:
        {
            "title": "Document Title",
            "sections": [
                {"heading": "Section 1", "content": "Paragraph text..."},
                {"heading": "Section 2", "content": "...", "bullet_points": ["Item 1", "Item 2"]}
            ]
        }
        """
        cls._ensure_dir()
        filename = f"report_{uuid.uuid4().hex[:8]}.docx"
        filepath = cls.BASE_DIR / filename
        
        doc = Document()
        
        # Title
        title = data.get("title", "Generated Report")
        doc.add_heading(title, 0)
        
        # Abstract/Intro
        if "abstract" in data:
            doc.add_paragraph(data["abstract"])
            
        # Sections
        sections = data.get("sections", [])
        for section in sections:
            if "heading" in section:
                doc.add_heading(section["heading"], level=1)
            
            if "content" in section:
                doc.add_paragraph(section["content"])
                
            if "bullet_points" in section and isinstance(section["bullet_points"], list):
                for point in section["bullet_points"]:
                    doc.add_paragraph(str(point), style='List Bullet')

        doc.save(str(filepath))
        
        return {
            "file_name": filename,
            "file_path": str(filepath),
            "download_url": cls._get_download_url(filename),
            "size_bytes": os.path.getsize(filepath)
        }

    @staticmethod
    def _normalize_xlsx_rows(rows: list) -> list:
        """
        容错处理：将各种异常格式的行数据统一转换为 [{"col1": v1, "col2": v2}, ...] 字典列表。

        处理的异常格式：
        1. 字符串列表 + tab分隔: ["col1\\tcol2", "v1\\tv2"] → 第一行做表头
        2. 字符串列表 + 竖线分隔: ["col1|col2", "v1|v2"]
        3. 字符串列表 + 逗号分隔: ["col1,col2", "v1,v2"]
        4. 嵌套列表: [["col1","col2"], ["v1","v2"]] → 第一行做表头
        5. 单值列表: ["a", "b", 1, 2] → 单列 {"值": x}
        """
        if not rows:
            return rows

        # 已经是标准的字典列表，直接返回
        if isinstance(rows[0], dict):
            return rows

        # 嵌套列表: [["col1","col2"], ["v1","v2"]]
        if isinstance(rows[0], (list, tuple)):
            headers = [str(h) for h in rows[0]]
            result = []
            for row in rows[1:]:
                record = {}
                for i, h in enumerate(headers):
                    record[h] = row[i] if i < len(row) else None
                result.append(record)
            return result if result else [{h: "" for h in headers}]

        # 字符串列表 — 尝试按分隔符拆列
        if isinstance(rows[0], str):
            # 检测分隔符：优先 tab > 竖线 > 逗号
            first_line = rows[0]
            sep = None
            if '\t' in first_line:
                sep = '\t'
            elif ' | ' in first_line:
                sep = ' | '
            elif '|' in first_line:
                sep = '|'

            if sep:
                headers = [h.strip() for h in first_line.split(sep)]
                result = []
                for line in rows[1:]:
                    vals = [v.strip() for v in line.split(sep)]
                    record = {}
                    for i, h in enumerate(headers):
                        record[h] = vals[i] if i < len(vals) else None
                    result.append(record)
                return result if result else [{h: "" for h in headers}]

            # 无分隔符，当作单列
            return [{"值": v} for v in rows]

        # 其他类型（数字等），当作单列
        return [{"值": v} for v in rows]

    @classmethod
    def generate_xlsx(cls, data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Generate Excel Spreadsheet
        Input data structure:
        {
            "sheets": {
                "Sheet1": [{"col1": 1, "col2": "A"}, ...],
                "Sheet2": [...]
            }
        }
        OR simple list for single sheet:
        {"data": [{"col1": 1, ...}]}
        """
        cls._ensure_dir()
        filename = f"data_{uuid.uuid4().hex[:8]}.xlsx"
        filepath = cls.BASE_DIR / filename
        
        # Determine data structure
        sheets_data = {}
        if "sheets" in data:
            sheets_data = data["sheets"]
        elif "data" in data:
            sheets_data = {"Sheet1": data["data"]}
        else:
            # Fallback for direct list input if passed incorrectly
            sheets_data = {"Sheet1": []}

        with pd.ExcelWriter(str(filepath), engine='openpyxl') as writer:
            for sheet_name, rows in sheets_data.items():
                if not rows:
                    pd.DataFrame({"Info": ["No Data"]}).to_excel(writer, sheet_name=sheet_name, index=False)
                else:
                    rows = cls._normalize_xlsx_rows(rows)
                    df = pd.DataFrame(rows)
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        return {
            "file_name": filename,
            "file_path": str(filepath),
            "download_url": cls._get_download_url(filename),
            "size_bytes": os.path.getsize(filepath)
        }

    @classmethod
    def generate_pptx(cls, data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Generate Professional PowerPoint Presentation
        Enhanced with:
        - Professional templates and layouts
        - Support for images and charts
        - Beautiful color schemes
        - Various slide types: title, content, quote, comparison, etc.

        Input data structure:
        {
            "title": "Presentation Title",
            "subtitle": "Subtitle",
            "theme": "business|gallery|beam|organic|celebration|circles|facet|integral",
            "slides": [
                {
                    "type": "title|content|quote|comparison|chart|image|section",
                    "title": "Slide Title",
                    "content": ["Bullet point 1", "Bullet point 2"] or "String",
                    "image_url": "base64 or URL",
                    "chart_data": {...},
                    "quote": "Quote text",
                    "author": "Quote author"
                }
            ]
        }
        """
        cls._ensure_dir()
        filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        filepath = cls.BASE_DIR / filename

        prs = cls._create_enhanced_presentation(data)

        prs.save(str(filepath))

        # Extract slide text for preview
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_content = []
            has_chart = False
            chart_type = ""

            # 首先检查原始数据，确定这个幻灯片是否应该是图表幻灯片
            original_slide_data = None
            if i < len(data.get("slides", [])):
                original_slide_data = data["slides"][i]
                if original_slide_data.get("type") == "chart":
                    has_chart = True
                    chart_type = original_slide_data.get("chart", {}).get("type", "")

            # 收集slide中的文本内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text = shape.text.strip()
                    slide_content.append(slide_text)

            # 根据原始数据添加图表类型标识
            if has_chart and chart_type:
                chart_type_text = {
                    "bar": "柱状图",
                    "pie": "饼图",
                    "line": "折线图",
                    "doughnut": "环形图",
                    "mindmap": "思维导图",
                    "timeline": "时间轴",
                    "column": "柱状图"
                }.get(chart_type, f"{chart_type}图表")

                # 使用不同的图标让前端更容易区分
                chart_icons = {
                    "mindmap": "🧠",
                    "timeline": "⏰"
                }
                icon = chart_icons.get(chart_type, "📊")
                slide_content.append(f"{icon} {chart_type_text}")

            # 检查图片内容
            for shape in slide.shapes:
                if type(shape).__name__ == "Picture":
                    if original_slide_data and original_slide_data.get("type") == "image":
                        slide_content.append("🖼️ 图片内容")

            slides_text.append("\n".join(slide_content))

        # Convert to base64 for frontend
        with open(filepath, "rb") as f:
            file_data = f.read()
        file_base64 = base64.b64encode(file_data).decode('utf-8')

        return {
            "file_name": filename,
            "file_path": str(filepath),
            "download_url": cls._get_download_url(filename),
            "file_size": os.path.getsize(filepath),
            "file_data_base64": file_base64,
            "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "slide_count": len(prs.slides),
            "slides_text": slides_text,
            "text_content": "\\n\\n".join([f"--- Slide {i+1} ---\\n{text}" for i, text in enumerate(slides_text)])
        }

    @classmethod
    def _create_enhanced_presentation(cls, data: Dict[str, Any]) -> Presentation:
        """Create an enhanced presentation with professional layouts"""
        prs = Presentation()

        # Set presentation properties
        prs.core_properties.title = data.get("title", "Professional Presentation")
        prs.core_properties.author = "AI Agent"

        # Color schemes
        themes = {
            "business": {"primary": RGBColor(0, 32, 96), "secondary": RGBColor(237, 125, 49), "accent": RGBColor(31, 78, 121)},
            "gallery": {"primary": RGBColor(91, 155, 213), "secondary": RGBColor(237, 125, 49), "accent": RGBColor(165, 165, 165)},
            "beam": {"primary": RGBColor(237, 137, 41), "secondary": RGBColor(165, 165, 165), "accent": RGBColor(91, 155, 213)},
            "organic": {"primary": RGBColor(134, 162, 53), "secondary": RGBColor(230, 230, 230), "accent": RGBColor(91, 155, 213)},
            "celebration": {"primary": RGBColor(237, 28, 36), "secondary": RGBColor(255, 217, 102), "accent": RGBColor(107, 107, 255)},
            "circles": {"primary": RGBColor(79, 129, 189), "secondary": RGBColor(247, 150, 70), "accent": RGBColor(189, 189, 189)},
            "facet": {"primary": RGBColor(158, 31, 99), "secondary": RGBColor(148, 138, 84), "accent": RGBColor(216, 202, 157)},
            "integral": {"primary": RGBColor(79, 129, 189), "secondary": RGBColor(192, 80, 77), "accent": RGBColor(155, 187, 89)}
        }

        theme = data.get("theme", "business")
        colors = themes.get(theme, themes["business"])

        # Title Slide
        cls._add_title_slide(prs, data, colors)

        # Process each slide
        for slide_data in data.get("slides", []):
            slide_type = slide_data.get("type", "content")
            title = slide_data.get("title", "Slide Title")

            # Check if slide contains chart data (even if type is not explicitly "chart")
            has_chart = bool(slide_data.get("chart") or slide_data.get("chart_type"))
            
            # Auto-detect type if missing based on content
            if slide_type == "content" or not slide_data.get("type"):
                if has_chart:
                    slide_type = "chart"
                elif slide_data.get("image_type") or slide_data.get("image_url"):
                    slide_type = "image"

            # 支持更多幻灯片类型
            if slide_type == "title":
                cls._add_title_only_slide(prs, title, colors)
            elif slide_type == "content" and not has_chart:
                cls._add_content_slide(prs, slide_data, colors)
            elif slide_type == "quote":
                cls._add_quote_slide(prs, slide_data, colors)
            elif slide_type == "comparison":
                cls._add_comparison_slide(prs, slide_data, colors)
            elif slide_type == "chart" or has_chart:
                cls._add_chart_slide(prs, slide_data, colors)
            elif slide_type == "image":
                cls._add_image_slide(prs, slide_data, colors)
            elif slide_type == "section":
                cls._add_section_slide(prs, title, colors)
            elif slide_type == "bullets":
                cls._add_bullets_slide(prs, slide_data, colors)
            elif slide_type == "timeline":
                cls._add_timeline_slide(prs, slide_data, colors)
            elif slide_type == "bullets_two_col":
                cls._add_two_column_slide(prs, slide_data, colors)
            else:
                cls._add_content_slide(prs, slide_data, colors)

        return prs

    @classmethod
    def _add_title_slide(cls, prs: Presentation, data: Dict[str, Any], colors: Dict):
        """Add professional title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add gradient background
        cls._add_gradient_background(slide, colors["primary"], colors["secondary"])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = data.get("title", "Professional Presentation")
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.add_paragraph()
        p.text = data.get("subtitle", "Generated with AI")
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(240, 240, 240)
        p.alignment = PP_ALIGN.CENTER

        # Add decorative shape
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(5), Inches(6), Inches(0.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.fill.background()

    @classmethod
    def _add_content_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add content slide with bullet points or formatted text"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide_data.get("title", "Content Slide")
        content = slide_data.get("content", slide_data.get("points", []))  # 支持 points 格式

        # Add header shape
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = colors["primary"]
        header_shape.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Add content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = content_box.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        # Backup for fallback to prevent blank slides
        raw_content_backup = str(content)

        if isinstance(content, str):
            content = [{"text": line.strip(), "level": 0 if not line.strip().startswith('  ') else 1} 
                      for line in content.split('\n') if line.strip()]
        elif isinstance(content, dict):
            # Flatten nested dicts (e.g. {'heading': '...', 'sectors': [...]})
            def _flatten_dict(d, lvl=0):
                res = []
                # Prioritize heading/title
                header = d.get("heading") or d.get("title") or d.get("name") or d.get("stage")
                if header: res.append({"text": str(header), "level": lvl})
                
                for k, v in d.items():
                    if k in ["heading", "title", "name", "stage", "chart", "image"]: continue
                    if isinstance(v, list):
                        for item in v:
                            res.extend(_flatten_dict(item, lvl + 1) if isinstance(item, dict) else [{"text": str(item), "level": lvl + 1}])
                    elif isinstance(v, dict):
                        res.extend(_flatten_dict(v, lvl + 1))
                    else:
                        # For simple values, maybe show key if it looks like a label
                        txt = f"{v}" # Just show value to keep it clean, or f"{k}: {v}"
                        if k in ["description", "event", "role"]: txt = str(v)
                        else: txt = str(v) 
                        res.append({"text": txt, "level": lvl + 1})
                return res
            content = _flatten_dict(content)
        
        # Final Fallback: If parsing resulted in empty content, use raw string
        if not content:
            content = [{"text": raw_content_backup if len(raw_content_backup) < 500 else raw_content_backup[:500] + "...", "level": 0}]

        # Ensure it's a list for the loop
        if not isinstance(content, list):
            content = [content]

        # Render content
        if isinstance(content, list):
            for item in content:
                p = tf.add_paragraph()
                if isinstance(item, dict):
                    # Try common keys for text content
                    text_content = item.get("text") or item.get("content") or item.get("value") or item.get("description") or str(item)
                    p.text = str(text_content)
                    p.level = int(item.get("level", 0))
                else:
                    p.text = str(item)
                p.font.size = Pt(20)
                p.space_after = Pt(12)

    @classmethod
    def _add_quote_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add elegant quote slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        quote = slide_data.get("quote", "")
        author = slide_data.get("author", "")

        # Add decorative left bar
        bar_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(0.3), Inches(5))
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = colors["accent"]

        # Add quote text
        quote_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(7), Inches(3))
        tf = quote_box.text_frame
        p = tf.add_paragraph()
        p.text = f'"{quote}"'
        p.font.size = Pt(36)
        p.font.italic = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # Add author
        if author:
            author_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(7), Inches(0.8))
            tf = author_box.text_frame
            p = tf.add_paragraph()
            p.text = f"— {author}"
            p.font.size = Pt(24)
            p.font.color.rgb = colors["secondary"]
            p.alignment = PP_ALIGN.RIGHT

    @classmethod
    def _add_comparison_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add side-by-side comparison slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title = slide_data.get("title", "Comparison")
        left_content = slide_data.get("left", [])
        right_content = slide_data.get("right", [])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # Add left column
        left_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(4.5))
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = RGBColor(240, 240, 240)

        left_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(3.6), Inches(4))
        tf = left_text.text_frame

        if isinstance(left_content, list):
            for item in left_content:
                p = tf.add_paragraph()
                p.text = str(item)
                p.font.size = Pt(18)
                p.space_after = Pt(10)

        # Add right column
        right_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2), Inches(4), Inches(4.5))
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(250, 250, 250)

        right_text = slide.shapes.add_textbox(Inches(5.7), Inches(2.2), Inches(3.6), Inches(4))
        tf = right_text.text_frame

        if isinstance(right_content, list):
            for item in right_content:
                p = tf.add_paragraph()
                p.text = str(item)
                p.font.size = Pt(18)
                p.space_after = Pt(10)

    @classmethod
    def _normalize_chart_data(cls, chart_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Centralized logic to extract standard labels/values from variable AI output formats.
        Mutates and returns chart_data with 'labels' and 'values' keys populated.
        """
        # 1. Check if already normalized
        if chart_data.get("labels") and chart_data.get("values"):
            return chart_data
        
        raw_data = chart_data.get("data")
        
        # 2. Handle List of Dictionaries (e.g. Timeline, sophisticated Bar data)
        # [{'year': 2020, 'val': 10}, {'year': 2021, 'val': 20}]
        if isinstance(raw_data, list) and raw_data and isinstance(raw_data[0], dict):
            # Try to find semantic keys
            keys = raw_data[0].keys()
            label_key = next((k for k in keys if k.lower() in ['year', 'years', 'date', 'month', 'name', 'category', 'label', 'region', 'model']), None)
            value_key = next((k for k in keys if k.lower() in ['value', 'values', 'count', 'amount', 'size', 'score', 'revenue', 'sales', 'market_size', 'percentage']), None)
            
            # Fallback if no specific keys found
            if not label_key: label_key = list(keys)[0]
            if not value_key and len(list(keys)) > 1: value_key = list(keys)[1]
            
            if label_key: 
                chart_data["labels"] = [str(item.get(label_key, "")) for item in raw_data]
            if value_key:
                # Type safe extraction
                vals = []
                for item in raw_data:
                    v = item.get(value_key, 0)
                    try:
                        vals.append(float(str(v).replace(',', '').replace('%', '')))
                    except:
                        vals.append(0)
                chart_data["values"] = vals
            return chart_data

        # 3. Handle Dictionary (Key-Value Map or Domain Keys)
        if isinstance(raw_data, dict):
            # 3a. Check for Domain Keys (Separate lists)
            # {'years': [2020, 2021], 'revenue': [100, 200]}
            found_labels = False
            found_values = False
            
            for k, v in raw_data.items():
                if not isinstance(v, list): continue
                k_lower = k.lower()
                
                # Label candidates
                if k_lower in ["years", "year", "month", "months", "date", "dates", "category", "categories", "names", "items", "models", "model", "labels"]:
                    chart_data["labels"] = v
                    found_labels = True
                
                # Value candidates
                elif k_lower in ["market_size", "size", "amount", "amounts", "value", "values", "count", "counts", "score", "scores", "revenue", "sales", "percentage", "share", "data"]:
                    chart_data["values"] = v
                    found_values = True
            
            if found_labels and found_values:
                return chart_data

            # 3b. Treat as direct Map {'Category': Value} if lists not found
            # Filter out non-data keys
            valid_map = {k: v for k, v in raw_data.items() if k not in ['type', 'title', 'subtitle', 'year']}
            if valid_map:
                chart_data["labels"] = list(valid_map.keys())
                vals = []
                for v in valid_map.values():
                    try:
                        vals.append(float(str(v).replace(',', '').replace('%', '')))
                    except:
                        vals.append(0)
                chart_data["values"] = vals
        
        return chart_data

    @classmethod
    def _add_chart_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add slide with chart (matplotlib generated)"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Robustly Extract Chart Object
        chart_obj = slide_data.get("chart", slide_data.get("chart_data", {}))
        
        # Handle "chart_type" + "data" flattening case
        if not chart_obj and "chart_type" in slide_data:
            chart_obj = {
                "type": slide_data.get("chart_type", "bar"),
                "data": slide_data.get("data", {}),
                "title": slide_data.get("title", "")
            }
            # Attempt to merge top-level data fields if data is missing or empty
            if not chart_obj["data"]:
                # Try to find loose data fields
                loose_data = {k:v for k,v in slide_data.items() if k in ['labels', 'values', 'categories', 'years', 'content']}
                if loose_data: chart_obj["data"] = loose_data

        chart_data = chart_obj
        title = chart_data.get("title", slide_data.get("title", "Chart Analysis"))
        chart_type = chart_data.get("type", chart_data.get("chart_type", "bar"))

        # --- NORMALIZE DATA ---
        chart_data = cls._normalize_chart_data(chart_data)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # 生成真正的图表
        try:
            print(f"[FileGenerator] Generating chart: type={chart_type}, data={chart_data}")
            # 重新配置中文字体
            configure_matplotlib_chinese()
            # 使用matplotlib生成图表
            img_buffer = io.BytesIO()

            if chart_type == "bar":
                # 支持两种数据结构：labels 或 categories
                labels = chart_data.get("labels", chart_data.get("categories", []))
                values = chart_data.get("values", [])
                chart_colors = chart_data.get("colors", [])

                print(f"[FileGenerator] Bar chart data: labels={labels}, values={values}, colors={chart_colors}")

                # 如果有data字段，从中提取数据
                data = chart_data.get("data", {})
                if not labels and data:
                    labels = data.get("labels", data.get("categories", []))
                    values = data.get("market_size", data.get("values", []))

                if labels and values and len(labels) == len(values):
                    print(f"[FileGenerator] Creating bar chart with {len(labels)} labels and {len(values)} values")
                    plt.figure(figsize=(10, 6))

                    # 使用指定的颜色或默认颜色
                    if chart_colors and len(chart_colors) >= len(labels):
                        bar_colors = [f"#{color}" for color in chart_colors]
                    else:
                        bar_colors = [colors["primary"].rgb if hasattr(colors["primary"], 'rgb') else "#1f77b4"] * len(labels)

                    bars = plt.bar(labels, values, color=bar_colors)
                    plt.title(chart_data.get("title", chart_data.get("chart_title", "数据分析")), fontsize=16)
                    plt.ylabel(chart_data.get("y_label", ""), fontsize=12)
                    plt.xticks(rotation=45)
                    plt.tight_layout()
                    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                    plt.close()

                    img_size = img_buffer.tell()
                    print(f"[FileGenerator] Chart generated: {img_size} bytes")
                    img_buffer.seek(0)

                    # 添加图表到幻灯片
                    slide.shapes.add_picture(img_buffer, Inches(1), Inches(2), Inches(8), Inches(4.5))
                    print("[FileGenerator] Chart added to slide successfully")
                else:
                    # 如果没有数据，使用文本表示
                    cls._add_chart_text_representation(slide, chart_data, colors)

            elif chart_type == "line":
                x_values = chart_data.get("x", chart_data.get("categories", chart_data.get("years", [])))
                y_values = chart_data.get("y", chart_data.get("values", []))
                chart_colors = chart_data.get("colors", [])

                print(f"[FileGenerator] Line chart data: x={x_values}, y={y_values}, colors={chart_colors}")

                if x_values and y_values and len(x_values) == len(y_values):
                    plt.figure(figsize=(10, 6))
                    line_color = f"#{chart_colors[0]}" if chart_colors else (colors["primary"].rgb if hasattr(colors["primary"], 'rgb') else "#1f77b4")

                    # 处理时间戳数据 - 确保x轴数据正确显示
                    if x_values and any(isinstance(x, (int, float)) and x > 1000000000 for x in x_values[:3]):
                        # 检测时间戳格式（Unix时间戳）
                        import datetime
                        try:
                            x_labels = [datetime.datetime.fromtimestamp(x).strftime('%Y-%m') if isinstance(x, (int, float)) and x > 1000000000 else str(x) for x in x_values]
                            plt.xticks(range(len(x_values)), x_labels, rotation=45)
                            print(f"[FileGenerator] Detected timestamp data, converted to: {x_labels[:3]}...")
                        except:
                            print("[FileGenerator] Failed to parse timestamps, using raw values")
                            plt.plot(x_values, y_values, color=line_color, linewidth=3, marker='o', markersize=8)
                    else:
                        # 常规数据
                        plt.plot(x_values, y_values, color=line_color, linewidth=3, marker='o', markersize=8)

                    plt.title(chart_data.get("title", chart_data.get("chart_title", "趋势分析")), fontsize=16)
                    plt.xlabel(chart_data.get("x_label", ""), fontsize=12)
                    plt.ylabel(chart_data.get("y_label", ""), fontsize=12)
                    plt.grid(True, alpha=0.3)
                    plt.tight_layout()
                    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                    plt.close()

                    img_buffer.seek(0)
                    slide.shapes.add_picture(img_buffer, Inches(1), Inches(2), Inches(8), Inches(4.5))
                else:
                    cls._add_chart_text_representation(slide, chart_data, colors)

            elif chart_type == "pie":
                labels = chart_data.get("labels", chart_data.get("categories", []))
                sizes = chart_data.get("values", [])
                chart_colors = chart_data.get("colors", [])

                print(f"[FileGenerator] Pie chart data: labels={labels}, sizes={sizes}, colors={chart_colors}")

                if labels and sizes and len(labels) == len(sizes):
                    plt.figure(figsize=(8, 8))

                    # 使用指定的颜色或默认颜色
                    if chart_colors and len(chart_colors) >= len(labels):
                        colors_for_pie = [f"#{color}" for color in chart_colors]
                    else:
                        colors_for_pie = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DDA0DD']

                    plt.pie(sizes, labels=labels, autopct='%1.1f%%', colors=colors_for_pie[:len(labels)])
                    plt.title(chart_data.get("title", chart_data.get("chart_title", "占比分析")), fontsize=16)
                    plt.axis('equal')
                    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                    plt.close()

                    img_buffer.seek(0)
                    slide.shapes.add_picture(img_buffer, Inches(2), Inches(1.5), Inches(6), Inches(6))
                else:
                    print(f"[FileGenerator] Invalid pie chart data: labels_len={len(labels) if labels else 0}, sizes_len={len(sizes) if sizes else 0}")
                    cls._add_chart_text_representation(slide, chart_data, colors)

            elif chart_type == "column":
                # Column chart is same as bar but with different orientation
                labels = chart_data.get("labels", chart_data.get("categories", []))
                values = chart_data.get("values", [])
                chart_colors = chart_data.get("colors", [])

                print(f"[FileGenerator] Column chart data: labels={labels}, values={values}, colors={chart_colors}")

                if labels and values and len(labels) == len(values):
                    plt.figure(figsize=(10, 6))

                    # 使用指定的颜色或默认颜色
                    if chart_colors and len(chart_colors) >= len(labels):
                        column_colors = [f"#{color}" for color in chart_colors]
                    else:
                        column_colors = [colors["primary"].rgb if hasattr(colors["primary"], 'rgb') else "#1f77b4"] * len(labels)

                    plt.bar(labels, values, color=column_colors)
                    plt.title(chart_data.get("title", chart_data.get("chart_title", "数据分析")), fontsize=16)
                    plt.ylabel(chart_data.get("y_label", ""), fontsize=12)
                    plt.xticks(rotation=45)
                    plt.tight_layout()
                    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                    plt.close()

                    img_size = img_buffer.tell()
                    print(f"[FileGenerator] Column chart generated: {img_size} bytes")
                    img_buffer.seek(0)

                    slide.shapes.add_picture(img_buffer, Inches(1), Inches(2), Inches(8), Inches(4.5))
                    print("[FileGenerator] Column chart added to slide successfully")
                else:
                    cls._add_chart_text_representation(slide, chart_data, colors)

            elif chart_type == "doughnut":
                labels = chart_data.get("labels", chart_data.get("categories", []))
                sizes = chart_data.get("values", [])
                chart_colors = chart_data.get("colors", [])

                print(f"[FileGenerator] Doughnut chart data: labels={labels}, sizes={sizes}, colors={chart_colors}")

                if labels and sizes and len(labels) == len(sizes):
                    plt.figure(figsize=(8, 8))

                    # 使用指定的颜色或默认颜色
                    if chart_colors and len(chart_colors) >= len(labels):
                        colors_for_doughnut = [f"#{color}" for color in chart_colors]
                    else:
                        colors_for_doughnut = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DDA0DD']

                    # 创建带中心空洞的饼图（环形图）
                    wedgeprops = {'width': 0.3, 'linewidth': 2, 'edgecolor': 'white'}
                    plt.pie(sizes, labels=labels, autopct='%1.1f%%', colors=colors_for_doughnut[:len(labels)],
                           wedgeprops=wedgeprops, pctdistance=0.85)
                    plt.title(chart_data.get("title", chart_data.get("chart_title", "占比分析")), fontsize=16)
                    plt.axis('equal')
                    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                    plt.close()

                    img_size = img_buffer.tell()
                    print(f"[FileGenerator] Doughnut chart generated: {img_size} bytes")
                    img_buffer.seek(0)

                    slide.shapes.add_picture(img_buffer, Inches(2), Inches(1.5), Inches(6), Inches(6))
                    print("[FileGenerator] Doughnut chart added to slide successfully")
                else:
                    cls._add_chart_text_representation(slide, chart_data, colors)

            elif chart_type == "mindmap" or chart_type == "mind":
                # 思维导图生成
                labels = chart_data.get("labels", chart_data.get("nodes", []))
                values = chart_data.get("values", chart_data.get("sizes", []))

                print(f"[FileGenerator] Mindmap data: labels={labels}, values={values}")

                if labels:
                    plt.figure(figsize=(12, 8))

                    # 创建径向布局的思维导图
                    import numpy as np

                    # 中心节点
                    center_label = chart_data.get("center", "核心主题")
                    center_size = chart_data.get("center_size", chart_data.get("values", [100])[0] if chart_data.get("values") else 100)

                    # 设置径向布局
                    n_nodes = len(labels)
                    if n_nodes > 0:
                        angles = np.linspace(0, 2*np.pi, n_nodes, endpoint=False)

                        # 设置颜色
                        chart_colors = chart_data.get("colors", [])
                        if chart_colors and len(chart_colors) >= n_nodes:
                            node_colors = [f"#{color}" for color in chart_colors[:n_nodes]]
                        else:
                            node_colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DDA0DD', '#FFA07A', '#98D8C8']

                        # 画中心节点
                        plt.scatter(0, 0, s=center_size*10, color=node_colors[0], alpha=0.8, edgecolors='white', linewidth=2)
                        plt.text(0, 0, center_label, ha='center', va='center', fontsize=12, fontweight='bold', color='white')

                        # 画子节点和连接线
                        for i, (label, angle) in enumerate(zip(labels, angles)):
                            # 子节点大小
                            node_size = values[i] if values and i < len(values) else 50

                            # 计算位置
                            x = np.cos(angle) * 3
                            y = np.sin(angle) * 3

                            # 画连接线
                            plt.plot([0, x], [0, y], color=node_colors[i % len(node_colors)], alpha=0.5, linewidth=2)

                            # 画子节点
                            plt.scatter(x, y, s=node_size*10, color=node_colors[i % len(node_colors)], alpha=0.8, edgecolors='white', linewidth=2)

                            # 添加标签
                            plt.text(x, y, str(label), ha='center', va='center', fontsize=10, fontweight='bold')

                        plt.title(chart_data.get("title", chart_data.get("chart_title", "思维导图")), fontsize=16)
                        plt.axis('off')
                        plt.tight_layout()
                        plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                        plt.close()

                        img_size = img_buffer.tell()
                        print(f"[FileGenerator] Mindmap generated: {img_size} bytes")
                        img_buffer.seek(0)

                        slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), Inches(8), Inches(6))
                        print("[FileGenerator] Mindmap added to slide successfully")
                    else:
                        cls._add_chart_text_representation(slide, chart_data, colors)
                else:
                    cls._add_chart_text_representation(slide, chart_data, colors)
            elif chart_type == "timeline":
                # 时间轴图表生成
                labels = chart_data.get("labels", chart_data.get("years", chart_data.get("dates", [])))
                values = chart_data.get("values", [])
                chart_colors = chart_data.get("colors", [])

                print(f"[FileGenerator] Timeline data: labels={labels}, values={values}")

                if labels and values and len(labels) == len(values):
                    plt.figure(figsize=(12, 6))

                    # 处理时间戳格式
                    import datetime
                    x_processed = []
                    for label in labels:
                        if isinstance(label, (int, float)) and label > 1000000000:
                            # Unix时间戳转换
                            try:
                                x_processed.append(datetime.datetime.fromtimestamp(label).strftime('%Y-%m'))
                            except:
                                x_processed.append(str(label))
                        else:
                            x_processed.append(str(label))

                    # 设置颜色
                    if chart_colors and len(chart_colors) >= len(labels):
                        bar_colors = [f"#{color}" for color in chart_colors[:len(labels)]]
                    else:
                        bar_colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7', '#DDA0DD']

                    # 创建横向柱状图作为时间轴
                    y_pos = np.arange(len(labels))
                    bars = plt.barh(y_pos, values, color=bar_colors, alpha=0.8)

                    # 添加数据标签
                    for i, (bar, value) in enumerate(zip(bars, values)):
                        plt.text(bar.get_width() + max(values)*0.01, bar.get_y() + bar.get_height()/2,
                               f'{value}', ha='left', va='center', fontsize=10)

                    plt.yticks(y_pos, x_processed)
                    plt.xlabel(chart_data.get("x_label", "数值"), fontsize=12)
                    plt.title(chart_data.get("title", chart_data.get("chart_title", "时间轴分析")), fontsize=16)
                    plt.grid(True, alpha=0.3, axis='x')
                    plt.tight_layout()
                    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                    plt.close()

                    img_size = img_buffer.tell()
                    print(f"[FileGenerator] Timeline generated: {img_size} bytes")
                    img_buffer.seek(0)

                    slide.shapes.add_picture(img_buffer, Inches(1), Inches(2), Inches(8), Inches(4.5))
                    print("[FileGenerator] Timeline added to slide successfully")
                else:
                    cls._add_chart_text_representation(slide, chart_data, colors)
            else:
                # 不支持的图表类型，使用文本表示
                print(f"[FileGenerator] Unsupported chart type: {chart_type}")
                cls._add_chart_text_representation(slide, chart_data, colors)

        except Exception as e:
            print(f"[FileGenerator] Chart generation failed: {e}")
            # 异常时使用文本表示
            cls._add_chart_text_representation(slide, chart_data, colors)

    @classmethod
    def _add_chart_text_representation(cls, slide, chart_data: Dict[str, Any], colors: Dict):
        """使用文本和形状创建图表表示"""
        chart_type = chart_data.get("type", "bar")
        chart_title = chart_data.get("title", chart_data.get("chart_title", "数据分析"))

        # 添加图表标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.8))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = f"📊 {chart_title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # 根据图表类型添加说明
        content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
        tf = content_box.text_frame

        if chart_type == "bar":
            p = tf.add_paragraph()
            p.text = "📈 柱状图分析"
            p.font.size = Pt(16)

            data = chart_data.get("data", {})
            if data:
                years = data.get("years", [])
                values = data.get("market_size", [])

                p = tf.add_paragraph()
                for i, (year, value) in enumerate(zip(years, values)):
                    p.add_run()
                    p.add_run().text = f"{year}: {value}亿美元"
                    p.level = 0
                    if i < len(years) - 1:
                        p.add_run().text = "\n"

        elif chart_type == "line":
            p = tf.add_paragraph()
            p.text = "📉 趋势图分析"
            p.font.size = Pt(16)

            p = tf.add_paragraph()
            p.text = "• 显示数据随时间的变化趋势"
            p = tf.add_paragraph()
            p.text = "• 支持多系列数据对比"
            p = tf.add_paragraph()
            p.text = "• 适用于时间序列分析"

        elif chart_type == "pie":
            p = tf.add_paragraph()
            p.text = "🥧 饼图分析"
            p.font.size = Pt(16)

            p = tf.add_paragraph()
            p.text = "• 显示各部分占比关系"
            p = tf.add_paragraph()
            p.text = "• 直观展示数据分布"
            p = tf.add_paragraph()
            p.text = "• 适用于组成结构分析"

        elif chart_type == "timeline":
            p = tf.add_paragraph()
            p.text = "🕒 时间轴"
            p.font.size = Pt(16)
            
            # 尝试解析时间轴数据
            events = chart_data.get("data", [])
            # 如果是归一化后的数据
            if not events and "labels" in chart_data:
                events = [{"year": l, "event": v} for l, v in zip(chart_data["labels"], chart_data["values"])]
            
            if isinstance(events, list):
                for item in events[:6]: # Limit to 6 items
                    if isinstance(item, dict):
                        year = item.get("year", item.get("time", ""))
                        event = item.get("event", item.get("description", ""))
                        if year or event:
                            p = tf.add_paragraph()
                            p.text = f"• {year}: {event}"
                            p.level = 0
                            p.font.size = Pt(14)

        else:
            p = tf.add_paragraph()
            p.text = f"📊 {chart_type.title()} 图表"
            p.font.size = Pt(16)

            p = tf.add_paragraph()
            p.text = "• 数据可视化展示"
            p = tf.add_paragraph()
            p.text = "• 请参考详细数据说明"

    @classmethod
    def _add_bullets_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add slide with bullet points"""
        cls._add_content_slide(prs, slide_data, colors)  # 复用内容幻灯片逻辑

    @classmethod
    def _add_timeline_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add slide with timeline content"""
        cls._add_content_slide(prs, slide_data, colors)  # 简化为内容幻灯片

    @classmethod
    def _add_two_column_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add slide with two column content"""
        cls._add_content_slide(prs, slide_data, colors)  # 简化为内容幻灯片

    @classmethod
    def _add_image_slide(cls, prs: Presentation, slide_data: Dict[str, Any], colors: Dict):
        """Add slide with image and caption"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title = slide_data.get("title", "Visual Content")
        image_url = slide_data.get("image_url", "")
        caption = slide_data.get("caption", "")

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        # Add placeholder for image (would need image processing)
        img_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), Inches(2),
            Inches(7), Inches(4)
        )
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = RGBColor(240, 240, 240)

        # Add image placeholder text
        img_text = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(4), Inches(1))
        tf = img_text.text_frame
        p = tf.add_paragraph()
        p.text = "[Image area]"
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(150, 150, 150)

        # Add caption
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.5))
            tf = caption_box.text_frame
            p = tf.add_paragraph()
            p.text = caption
            p.font.size = Pt(16)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = colors["secondary"]

    @classmethod
    def _add_section_slide(cls, prs: Presentation, title: str, colors: Dict):
        """Add section divider slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Add full background
        background_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
        background_shape.fill.solid()
        background_shape.fill.fore_color.rgb = colors["primary"]

        # Add section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title.upper()
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    @classmethod
    def _add_title_only_slide(cls, prs: Presentation, title: str, colors: Dict):
        """Add title only slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Add title with design
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.CENTER

    @classmethod
    def _add_gradient_background(cls, slide, color1: RGBColor, color2: RGBColor):
        """Add gradient background to slide"""
        # Placeholder for gradient background
        # Full implementation would require XML manipulation
        pass

    @classmethod
    def _get_chart_colors(cls, colors: Dict) -> List[RGBColor]:
        """Get color palette for charts"""
        return [
            colors.get("primary", RGBColor(79, 129, 189)),
            colors.get("secondary", RGBColor(192, 80, 77)),
            colors.get("accent", RGBColor(155, 187, 89)),
            RGBColor(255, 127, 0),
            RGBColor(128, 128, 128)
        ]

    @classmethod
    def generate_pdf(cls, data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Generate PDF Document with Chinese Support and Markdown Parsing
        """
        print(f"[FileGenerator] Generating PDF. Data type: {type(data)}")
        if data is None:
            print("[FileGenerator] Error: data is None")
            raise ValueError("Data cannot be None")
        if not isinstance(data, dict):
             print(f"[FileGenerator] Error: data is not a dict, it is {type(data)}")
             # Try to recover if it's a string that looks like JSON? Or just wrap it?
             # For now, just wrap it as content if it's a string
             if isinstance(data, str):
                 data = {"content": data, "title": "Generated Document"}
             else:
                 raise ValueError(f"Data must be a dictionary, got {type(data)}")

        cls._ensure_dir()
        filename = f"document_{uuid.uuid4().hex[:8]}.pdf"
        filepath = cls.BASE_DIR / filename
        
        # 1. Register Chinese Fonts
        font_name = "Helvetica" # Default
        try:
            # Common Unicode fonts
            fallback_fonts = [
               ("ZenHei", "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"), # Prioritize this as we will install it
               ("SimHei", "SimHei.ttf"),
               ("Microsoft YaHei", "msyh.ttc"),
               ("Heiti SC", "/System/Library/Fonts/STHeiti Light.ttc"),
               ("PingFang SC", "/System/Library/Fonts/PingFang.ttc"),
               ("Noto Sans CJK SC", "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
               ("Arial Unicode MS", "Arial Unicode.ttf")
            ]
            
            for name, path in fallback_fonts:
                try:
                     # Check if file exists or is builtin
                     if os.path.exists(path) or path.endswith(".ttf") or path.endswith(".ttc"):
                        if path.endswith(".ttc"):
                             from reportlab.pdfbase.ttfonts import TTFont
                             pdfmetrics.registerFont(TTFont(name, path, subfontIndex=0))
                        else:
                             from reportlab.pdfbase.ttfonts import TTFont
                             pdfmetrics.registerFont(TTFont(name, path))
                        font_name = name
                        print(f"[FileGenerator] Registered PDF font: {name}")
                        break
                except Exception as e:
                    # Ignore registration errors
                    continue
        except Exception as e:
            print(f"[FileGenerator] Font registration failed: {e}")

        # 2. Build PDF using Platypus (Better for wrapping)
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
        from reportlab.lib.units import inch, cm
        from reportlab.lib import colors
        import re

        doc = SimpleDocTemplate(
            str(filepath), 
            pagesize=A4,
            rightMargin=72, leftMargin=72,
            topMargin=72, bottomMargin=72
        )
        story = []
        
        styles = getSampleStyleSheet()
        
        # --- Custom Styles ---
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontName=font_name, # Use the registered Chinese font
            fontSize=26,
            leading=32,
            spaceAfter=20,
            alignment=TA_CENTER,
            textColor=colors.HexColor('#1a5276')
        )
        
        # Headings - Make them VERY distinct
        h1_style = ParagraphStyle(
            'CustomH1',
            parent=styles['Normal'],
            fontName=font_name,
            fontSize=20,
            leading=26,
            spaceBefore=18,
            spaceAfter=10,
            textColor=colors.HexColor('#1f618d'), # Dark Blue
            keepWithNext=True
        )
        
        h2_style = ParagraphStyle(
            'CustomH2',
            parent=styles['Normal'],
            fontName=font_name,
            fontSize=16,
            leading=22,
            spaceBefore=14,
            spaceAfter=8,
            textColor=colors.HexColor('#2471a3'), # Blue
            keepWithNext=True
        )

        h3_style = ParagraphStyle(
            'CustomH3',
            parent=styles['Normal'],
            fontName=font_name,
            fontSize=14,
            leading=18,
            spaceBefore=10,
            spaceAfter=6,
            textColor=colors.HexColor('#5499c7'), # Light Blue
            keepWithNext=True
        )
        
        # Body Text
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontName=font_name,
            fontSize=11,
            leading=16,
            spaceAfter=6,
            alignment=TA_JUSTIFY,
            textColor=colors.black
        )

        # Bullet List Item Style
        bullet_style = ParagraphStyle(
            'CustomBullet',
            parent=body_style,
            leftIndent=0,   # ListFlowable handles indentation
            firstLineIndent=0,
            spaceAfter=0
        )
        
        # Add Title (if provided in data)
        doc_title = data.get("title", "")
        if doc_title:
            story.append(Paragraph(doc_title, title_style))
            story.append(Spacer(1, 10))
            story.append(Paragraph("<br/><br/>", body_style))

        # Add Content with Markdown Parsing
        raw_content = data.get("content", "")
        if raw_content is None: raw_content = ""
        
        # --- PRE-PROCESSING & CLEANUP ---
        # 1. Handle literal newlines
        if "\\n" in raw_content:
             raw_content = raw_content.replace("\\n", "\n")
        
        # 2. Fix "Glued" content using REGEX
        # Insert newline before bullets (●, •) if safe
        raw_content = re.sub(r'([^\n])\s*([●•])', r'\1\n\2', raw_content)
        
        # Insert newline before Numbered Lists (1., 2.)
        # Look for " 1. " or ". 1. " pattern
        raw_content = re.sub(r'(\s)(\d+\.)', r'\n\2', raw_content)

        # Insert newline after Question Marks
        raw_content = re.sub(r'([？?])(\s*[^\n])', r'\1\n\2', raw_content)
        
        # Replace "Fullwidth" characters ensuring consistent parsing
        raw_content = raw_content.replace("．", ".") 
        
        print(f"[FileGenerator] Content length after preprocessing: {len(raw_content)}")

        # Helper to process bold text: **text** -> <b>text</b>
        def process_text_style(text):
            # Handle Bold **text**
            parts = text.split('**')
            new_text = ""
            for i, part in enumerate(parts):
                if i % 2 == 1: new_text += f"<b>{part}</b>"
                else: new_text += part
            text = new_text
            
            # Handle Bold __text__
            if '__' in text:
                parts = text.split('__')
                new_text = ""
                for i, part in enumerate(parts):
                    if i % 2 == 1: new_text += f"<b>{part}</b>"
                    else: new_text += part
                text = new_text
                
            return text

        lines = [line.strip() for line in raw_content.split('\n') if line.strip()]
        current_list_items = []

        def flush_list(items, s):
            if items:
                # Create a ListFlowable
                list_flow = ListFlowable(
                    [ListItem(Paragraph(process_text_style(item), bullet_style)) for item in items],
                    bulletType='bullet',
                    start='circle',
                    leftIndent=20,
                    bulletFontSize=8,
                    spaceAfter=6
                )
                s.append(list_flow)
                s.append(Spacer(1, 4))
                items.clear()

        for line in lines:
            # Skip if it's an exact duplicate of the title (common repetitive artifact)
            if doc_title and line == doc_title:
                continue
                
            line = line.strip()
            if not line:
                flush_list(current_list_items, story)
                continue

            # Parse Markdown Headers
            if line.startswith('# '):
                flush_list(current_list_items, story)
                text = line[2:].strip()
                story.append(Paragraph(process_text_style(text), h1_style))
            
            elif line.startswith('## '):
                flush_list(current_list_items, story)
                text = line[3:].strip()
                story.append(Paragraph(process_text_style(text), h2_style))
                
            elif line.startswith('### '):
                flush_list(current_list_items, story)
                text = line[4:].strip()
                story.append(Paragraph(process_text_style(text), h3_style))
                
            # Parse Bullet Points (Standard -, *, and now ●, •)
            elif line.startswith('- ') or line.startswith('* ') or line.startswith('●') or line.startswith('•'):
                # Determine where text starts
                if line.startswith('- ') or line.startswith('* '):
                    text = line[2:].strip()
                else: 
                     # For ● and •, usually just 1 char
                     text = line[1:].strip()
                current_list_items.append(text)
                
            # Parse Numbered Lists (1. Item)
            elif len(line) > 2 and line[0].isdigit() and '. ' in line[:5]:
                 flush_list(current_list_items, story)
                 dot_index = line.find('. ')
                 text = line[dot_index+2:].strip()
                 number = line[:dot_index+1]
                 p = Paragraph(f"<b>{number}</b> {process_text_style(text)}", body_style)
                 story.append(p)
            
            # Normal Text
            else:
                flush_list(current_list_items, story)
                story.append(Paragraph(process_text_style(line), body_style))
        
        # Flush any remaining list items
        flush_list(current_list_items, story)

        doc.build(story)

        return {
            "file_name": filename,
            "file_path": str(filepath),
            "download_url": cls._get_download_url(filename),
            "size_bytes": os.path.getsize(filepath)
        }


# ========== 远程文件下载工具函数 ==========

def download_remote_file(file_url: str, allowed_extensions: list = None) -> dict:
    """
    下载远程文件（HTTP URL / MinIO URL）到本地 data/generated/ 目录，
    返回本地路径和可访问的 download_url。

    Args:
        file_url: 远程文件 URL
        allowed_extensions: 允许的文件扩展名列表，如 ['.xlsx', '.xls']

    Returns:
        {"file_path": str, "download_url": str, "file_name": str}

    Raises:
        ValueError: URL 无效或文件类型不支持
        RuntimeError: 下载失败
    """
    import requests
    import uuid
    from urllib.parse import urlparse, unquote

    if not file_url or not file_url.startswith(("http://", "https://")):
        raise ValueError(f"无效的文件 URL: {file_url}")

    # 从 URL 中提取文件名
    parsed = urlparse(file_url)
    url_path = unquote(parsed.path)
    original_name = url_path.split("/")[-1] if "/" in url_path else "downloaded_file"

    # 检查扩展名
    ext = ""
    if "." in original_name:
        ext = "." + original_name.rsplit(".", 1)[-1].lower()

    if allowed_extensions and ext not in allowed_extensions:
        raise ValueError(f"不支持的文件类型: {ext}，允许: {allowed_extensions}")

    # 生成本地文件名（保留原名 + UUID 防冲突）
    short_id = str(uuid.uuid4())[:8]
    if ext:
        base_name = original_name.rsplit(".", 1)[0]
        local_name = f"{base_name}_{short_id}{ext}"
    else:
        local_name = f"{original_name}_{short_id}"

    # 确保目录存在
    save_dir = get_generated_files_dir()
    save_dir.mkdir(parents=True, exist_ok=True)
    local_path = save_dir / local_name

    # 下载文件
    try:
        resp = requests.get(file_url, timeout=60, stream=True)
        resp.raise_for_status()

        with open(local_path, "wb") as f:
            for chunk in resp.iter_content(chunk_size=8192):
                f.write(chunk)

    except Exception as e:
        raise RuntimeError(f"文件下载失败: {file_url} -> {e}")

    # 生成可访问的下载 URL
    download_url = FileGenerator._get_download_url(local_name)

    return {
        "file_path": str(local_path),
        "download_url": download_url,
        "file_name": local_name,
    }


def ensure_download_url(file_path: str, download_url: str = "") -> str:
    """
    确保文件有可访问的 download_url。
    如果已有 download_url 则直接返回；
    否则将文件上传到 MinIO 返回公网 URL。

    Args:
        file_path: 本地文件路径
        download_url: 已有的下载地址（可能为空）

    Returns:
        可访问的下载 URL
    """
    if download_url:
        return download_url

    # 文件不存在则返回空
    if not file_path or not os.path.exists(file_path):
        return ""

    file_name = os.path.basename(file_path)

    generated_dir = get_generated_files_dir().resolve()
    # 如果文件在生成目录下，直接生成 agent-service 下载 URL
    if generated_dir == Path(file_path).resolve() or str(Path(file_path).resolve()).startswith(f"{generated_dir}{os.sep}"):
        return FileGenerator._get_download_url(file_name)

    # 其他路径的文件，统一暴露为本地下载地址
    try:
        return build_download_url_for_path(file_path)
    except Exception:
        return ""

# ========== 本地对象存储工具函数 ==========

def upload_file_to_local_storage(
    file_data: bytes,
    user_id: int,
    file_name: str,
    content_type: str = "application/octet-stream",
) -> str:
    """
    将文件保存到本地对象存储，并返回下载 URL。

    Args:
        file_data: 文件二进制内容
        user_id: 用户 ID（用于路径隔离）
        file_name: 原始文件名
        content_type: MIME 类型

    Returns:
        本地文件下载 URL，失败返回空字符串
    """
    from loguru import logger

    try:
        storage = save_bytes_to_local_storage(
            file_data=file_data,
            bucket_name="uploads",
            object_name=f"{user_id}/{file_name}",
        )
        logger.info(
            f"[LocalUpload] Stored: uploads/{storage['object_name']} "
            f"({len(file_data)} bytes, content_type={content_type})"
        )
        return storage["download_url"]
    except Exception as e:
        logger.error(f"[LocalUpload] Failed: {e}")
        return ""
