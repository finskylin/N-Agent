"""
通用文档解析工具层

从 document_reader skill 抽取的纯工具函数。
skill 和 knowledge_file_service 各自引用此模块，互不依赖。

支持格式: PDF / DOCX / XLSX / MD / TXT / JSON / CSV / 图片(OCR)
"""

import os
import io
import base64
import json as _json
from typing import Dict, Any, Optional
from pathlib import Path

from loguru import logger


# --- 文件类型映射（从配置读取时可覆盖） ---

EXT_TO_TYPE = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".md": "md",
    ".txt": "txt",
    ".json": "json",
    ".csv": "csv",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".bmp": "image",
    ".webp": "image",
}

EXT_TO_MIME = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls": "application/vnd.ms-excel",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".md": "text/markdown",
    ".txt": "text/plain",
    ".json": "application/json",
    ".csv": "text/csv",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
    ".webp": "image/webp",
}

OCR_MIN_CHARS = 50


def detect_file_type(file_path: str) -> Optional[str]:
    """根据扩展名检测文件类型"""
    ext = Path(file_path).suffix.lower()
    return EXT_TO_TYPE.get(ext)


def detect_mime_type(file_path: str) -> str:
    """根据扩展名检测 MIME 类型"""
    ext = Path(file_path).suffix.lower()
    return EXT_TO_MIME.get(ext, "application/octet-stream")


async def parse_file(file_path: str, max_chars: int = 100000) -> Dict[str, Any]:
    """
    统一入口：根据文件类型调用对应解析器

    Args:
        file_path: 文件绝对路径
        max_chars: 最大字符数限制

    Returns:
        解析结果字典，至少包含:
            file_name, file_type, file_size, text_content
    """
    file_type = detect_file_type(file_path)
    if not file_type:
        return {
            "file_name": os.path.basename(file_path),
            "file_type": "unknown",
            "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
            "text_content": "",
            "error": f"Unsupported file extension: {Path(file_path).suffix}",
        }

    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "xlsx": parse_xlsx,
        "pptx": parse_pptx,
        "md": parse_text,
        "txt": parse_text,
        "json": parse_json,
        "csv": parse_csv,
        "image": parse_image,
    }

    parser = parsers.get(file_type)
    if not parser:
        return {
            "file_name": os.path.basename(file_path),
            "file_type": file_type,
            "file_size": os.path.getsize(file_path),
            "text_content": "",
            "error": f"No parser for type: {file_type}",
        }

    try:
        return await parser(file_path, max_chars)
    except Exception as e:
        logger.error(f"[DocumentParser] Error parsing {file_path}: {e}")
        return {
            "file_name": os.path.basename(file_path),
            "file_type": file_type,
            "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
            "text_content": "",
            "error": str(e),
        }


# ============================================================
# 各类型解析器
# ============================================================

async def parse_pdf(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析 PDF（pdfplumber 优先，pypdf 兜底）"""
    text_parts = []
    tables = []
    page_count = 0
    metadata = {}

    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            metadata = pdf.metadata or {}

            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    text_parts.append(text.strip())

                page_tables = page.extract_tables()
                for tbl in page_tables:
                    if tbl and len(tbl) > 1:
                        headers = [str(h or "") for h in tbl[0]]
                        rows = [[str(c or "") for c in row] for row in tbl[1:]]
                        tables.append({
                            "location": f"Page {i + 1}",
                            "headers": headers,
                            "rows": rows[:50],
                        })

                if len("\n".join(text_parts)) >= max_chars:
                    break

    except Exception as e:
        logger.warning(f"[DocumentParser] pdfplumber failed, trying pypdf: {e}")
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            page_count = len(reader.pages)
            metadata = dict(reader.metadata) if reader.metadata else {}

            for page in reader.pages:
                text = page.extract_text()
                if text and text.strip():
                    text_parts.append(text.strip())
                if len("\n".join(text_parts)) >= max_chars:
                    break
        except Exception as e2:
            logger.error(f"[DocumentParser] pypdf also failed: {e2}")

    full_text = "\n\n".join(text_parts)
    return {
        "file_name": os.path.basename(file_path),
        "file_type": "pdf",
        "file_size": os.path.getsize(file_path),
        "page_count": page_count,
        "text_content": full_text[:max_chars],
        "tables": tables[:20],
        "word_count": len(full_text),
        "metadata": {k: str(v) for k, v in metadata.items() if v} if metadata else {},
    }


async def parse_docx(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析 DOCX"""
    from docx import Document

    doc = Document(file_path)
    text_parts = []
    tables = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
        if len("\n".join(text_parts)) >= max_chars:
            break

    for i, tbl in enumerate(doc.tables):
        rows_data = []
        headers = []
        for j, row in enumerate(tbl.rows):
            cells = [cell.text.strip() for cell in row.cells]
            if j == 0:
                headers = cells
            else:
                rows_data.append(cells)
        if headers:
            tables.append({
                "location": f"Table {i + 1}",
                "headers": headers,
                "rows": rows_data[:50],
            })

    full_text = "\n".join(text_parts)
    return {
        "file_name": os.path.basename(file_path),
        "file_type": "docx",
        "file_size": os.path.getsize(file_path),
        "text_content": full_text[:max_chars],
        "tables": tables[:20],
        "word_count": len(full_text),
        "metadata": {
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
        },
    }


async def parse_xlsx(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析 XLSX"""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    sheets_info = []
    all_text_parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))

        if not rows:
            sheets_info.append({"sheet": sheet_name, "rows": 0, "cols": 0})
            continue

        headers = [str(h or "") for h in rows[0]]
        data_rows = [[str(c or "") for c in row] for row in rows[1:]]

        sheets_info.append({
            "sheet": sheet_name,
            "rows": len(data_rows),
            "cols": len(headers),
            "headers": headers,
            "sample_rows": data_rows[:30],
        })

        text_repr = f"--- Sheet: {sheet_name} ---\n"
        text_repr += " | ".join(headers) + "\n"
        for row in data_rows[:30]:
            text_repr += " | ".join(row) + "\n"
        all_text_parts.append(text_repr)

    wb.close()
    full_text = "\n\n".join(all_text_parts)

    return {
        "file_name": os.path.basename(file_path),
        "file_type": "xlsx",
        "file_size": os.path.getsize(file_path),
        "text_content": full_text[:max_chars],
        "tables": sheets_info,
        "word_count": len(full_text),
        "metadata": {"sheet_count": len(sheets_info)},
    }


async def parse_pptx(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析 PPTX"""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []
    slide_count = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
        if len("\n".join(text_parts)) >= max_chars:
            break

    full_text = "\n\n".join(text_parts)
    return {
        "file_name": os.path.basename(file_path),
        "file_type": "pptx",
        "file_size": os.path.getsize(file_path),
        "page_count": slide_count,
        "text_content": full_text[:max_chars],
        "word_count": len(full_text),
        "metadata": {"slide_count": slide_count},
    }


async def parse_text(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析纯文本（MD / TXT）"""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read(max_chars)

    return {
        "file_name": os.path.basename(file_path),
        "file_type": Path(file_path).suffix.lstrip(".") or "txt",
        "file_size": os.path.getsize(file_path),
        "text_content": content,
        "word_count": len(content),
        "metadata": {},
    }


async def parse_json(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析 JSON"""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        raw = f.read(max_chars * 2)

    try:
        data = _json.loads(raw)
        formatted = _json.dumps(data, ensure_ascii=False, indent=2)
    except _json.JSONDecodeError:
        formatted = raw

    return {
        "file_name": os.path.basename(file_path),
        "file_type": "json",
        "file_size": os.path.getsize(file_path),
        "text_content": formatted[:max_chars],
        "word_count": len(formatted),
        "metadata": {},
    }


async def parse_csv(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析 CSV"""
    import csv

    text_parts = []
    row_count = 0
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            text_parts.append(" | ".join(row))
            row_count += 1
            if len("\n".join(text_parts)) >= max_chars:
                break

    full_text = "\n".join(text_parts)
    return {
        "file_name": os.path.basename(file_path),
        "file_type": "csv",
        "file_size": os.path.getsize(file_path),
        "text_content": full_text[:max_chars],
        "word_count": len(full_text),
        "metadata": {"row_count": row_count},
    }


async def parse_image(file_path: str, max_chars: int) -> Dict[str, Any]:
    """解析图片：pytesseract OCR 优先，视觉模型兜底"""
    ocr_text = ""

    try:
        import pytesseract
        from PIL import Image

        img = Image.open(file_path)
        ocr_text = pytesseract.image_to_string(img, lang="chi_sim+eng").strip()
        logger.info(f"[DocumentParser] pytesseract OCR: {len(ocr_text)} chars")
    except ImportError:
        logger.warning("[DocumentParser] pytesseract not installed, skipping OCR")
    except Exception as e:
        logger.warning(f"[DocumentParser] pytesseract failed: {e}")

    if len(ocr_text) >= OCR_MIN_CHARS:
        return {
            "file_name": os.path.basename(file_path),
            "file_type": "image",
            "file_size": os.path.getsize(file_path),
            "text_content": ocr_text[:max_chars],
            "word_count": len(ocr_text),
            "metadata": {"recognition_method": "ocr"},
        }

    # OCR 不足，尝试视觉模型
    vision_text = await _call_vision_model(file_path)
    text_content = vision_text or ocr_text or ""

    return {
        "file_name": os.path.basename(file_path),
        "file_type": "image",
        "file_size": os.path.getsize(file_path),
        "text_content": text_content[:max_chars],
        "word_count": len(text_content),
        "metadata": {
            "recognition_method": "vision_model" if vision_text else "ocr",
        },
    }


async def _call_vision_model(file_path: str) -> str:
    """调用视觉大模型识别图片"""
    try:
        import httpx
    except ImportError:
        return ""

    from app.config import settings
    api_key = settings.vision_api_key or settings.anthropic_auth_token or ""
    base_url = settings.vision_api_base_url or settings.anthropic_base_url or ""

    if not api_key or not base_url:
        return ""

    with open(file_path, "rb") as f:
        image_data = base64.b64encode(f.read()).decode("utf-8")

    ext = Path(file_path).suffix.lower()
    media_type = EXT_TO_MIME.get(ext, "image/png")

    vision_model = settings.vision_model

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                f"{base_url}/chat/completions",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": vision_model,
                    "messages": [{
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "请详细识别并描述这张图片中的所有文字内容和关键信息。如果包含表格，请以结构化方式输出。",
                            },
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{media_type};base64,{image_data}"},
                            },
                        ],
                    }],
                    "max_tokens": 4096,
                },
            )
            result = resp.json()
        return result.get("choices", [{}])[0].get("message", {}).get("content", "")
    except Exception as e:
        logger.error(f"[DocumentParser] Vision model failed: {e}")
        return ""
