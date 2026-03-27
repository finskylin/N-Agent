#!/usr/bin/env python3
"""创建测试用的 Office 文件"""
import os
from docx import Document
from openpyxl import Workbook
from fpdf import FPDF
from pptx import Presentation

# 创建目录
test_dir = "/Users/finlin/ai_workspace/master_growth/output_project/sthg_agent_service/test_files"
os.makedirs(test_dir, exist_ok=True)

# 1. 创建 Word 文档
doc = Document()
doc.add_heading('股票分析报告', 0)
doc.add_paragraph('这是一份测试用的股票分析报告文档。')
doc.add_paragraph('本报告包含以下内容：')
doc.add_paragraph('1. 市场概况分析', style='List Number')
doc.add_paragraph('2. 技术指标分析', style='List Number')
doc.add_paragraph('3. 投资建议', style='List Number')
doc.add_paragraph('')

doc.add_heading('市场概况', 1)
doc.add_paragraph('当前市场呈现震荡上涨趋势，成交量稳步放大。')

doc.add_heading('技术指标', 1)
doc.add_paragraph('MACD 指标显示金叉信号，KDJ 指标进入强势区域。')

doc.add_heading('投资建议', 1)
doc.add_paragraph('建议投资者关注该股票的后续走势，逢低吸纳。')

docx_path = os.path.join(test_dir, "股票分析报告.docx")
doc.save(docx_path)
print(f"✅ 创建 Word 文档: {docx_path}")

# 2. 创建 Excel 表格
wb = Workbook()
ws = wb.active
ws.title = "股票数据"

# 添加表头
headers = ["日期", "开盘价", "最高价", "最低价", "收盘价", "成交量"]
for col, header in enumerate(headers, 1):
    ws.cell(row=1, column=col, value=header)

# 添加示例数据
data = [
    ["2024-01-15", 10.50, 10.80, 10.40, 10.75, 125000],
    ["2024-01-16", 10.70, 11.00, 10.65, 10.95, 138000],
    ["2024-01-17", 10.90, 11.20, 10.85, 11.15, 142000],
    ["2024-01-18", 11.10, 11.40, 11.05, 11.35, 156000],
    ["2024-01-19", 11.30, 11.60, 11.25, 11.55, 168000],
]

for row_idx, row_data in enumerate(data, 2):
    for col_idx, value in enumerate(row_data, 1):
        ws.cell(row=row_idx, column=col_idx, value=value)

xlsx_path = os.path.join(test_dir, "股票数据.xlsx")
wb.save(xlsx_path)
print(f"✅ 创建 Excel 表格: {xlsx_path}")

# 3. 创建 PDF 文档
pdf_path = os.path.join(test_dir, "技术分析报告.pdf")
pdf = FPDF()
pdf.add_page()

# 设置中文字体（使用系统自带字体）
pdf.set_font("Arial", 'B', 18)
pdf.cell(0, 10, "Technical Analysis Report", ln=True, align='C')

pdf.set_font("Arial", '', 12)
pdf.ln(10)
pdf.cell(0, 10, "This report analyzes stock technical indicators.", ln=True)
pdf.ln(5)
pdf.cell(0, 10, "Key Findings:", ln=True)
pdf.ln(3)
pdf.cell(10, 8, "- MACD golden cross formed", ln=True)
pdf.cell(10, 8, "- RSI at 65, in strong zone", ln=True)
pdf.cell(10, 8, "- Volume continues to increase", ln=True)
pdf.cell(10, 8, "- Moving average system bullish", ln=True)

pdf.output(pdf_path)
print(f"✅ 创建 PDF 文档: {pdf_path}")

# 4. 创建 PowerPoint 演示文稿
prs = Presentation()

# 幻灯片 1: 标题
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
title.text = "Stock Investment Analysis"
subtitle = slide1.placeholders[1]
subtitle.text = "Q1 2024 Report"

# 幻灯片 2: 市场分析
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide2.shapes
title_shape = shapes.title
title_shape.text = "Market Analysis"

body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Key Points"

p = tf.add_paragraph()
p.text = "Overall market trend upward"
p.level = 1
p = tf.add_paragraph()
p.text = "Sector rotation evident"
p.level = 1
p = tf.add_paragraph()
p.text = "Continuous fund inflow"
p.level = 1

# 幻灯片 3: 投资建议
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide3.shapes
title_shape = shapes.title
title_shape.text = "Investment Advice"

body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Strategy"

p = tf.add_paragraph()
p.text = "Buy on dips, control position"
p.level = 1
p = tf.add_paragraph()
p.text = "Monitor technical support levels"
p.level = 1
p = tf.add_paragraph()
p.text = "Set reasonable stop-loss points"
p.level = 1

pptx_path = os.path.join(test_dir, "投资分析演示.pptx")
prs.save(pptx_path)
print(f"✅ 创建 PowerPoint 演示文稿: {pptx_path}")

print("\n" + "="*50)
print("所有测试文件创建完成！")
print("="*50)
print(f"测试文件目录: {test_dir}")
print("- " * 30)
for f in [docx_path, xlsx_path, pdf_path, pptx_path]:
    size = os.path.getsize(f)
    print(f"{os.path.basename(f)}: {size} bytes")
