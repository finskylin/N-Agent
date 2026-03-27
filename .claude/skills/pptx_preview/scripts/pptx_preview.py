"""
PPTX Preview Skill
PowerPoint 演示文稿预览技能 - 将 pptx 文件渲染到数据看板
"""

# Imports

# ── file_storage helpers (no app.services dependency) ────────────────────────
from typing import Any, Dict, List
import os as _os
import uuid as _uuid
from pathlib import Path as _Path
from urllib.parse import quote as _quote, urlparse as _urlparse, unquote as _unquote

def _get_public_base_url() -> str:
    base = _os.getenv("AGENT_PUBLIC_BASE_URL", "").rstrip("/")
    if base:
        return base
    host = _os.getenv("AGENT_EXTERNAL_HOST", "127.0.0.1")
    port = _os.getenv("AGENT_SERVICE_PORT", "8000")
    return f"http://{host}:{port}"

def _get_generated_files_dir() -> _Path:
    configured = _os.getenv("LOCAL_GENERATED_FILES_DIR", "").strip()
    return _Path(configured) if configured else _Path("app/data/generated")

def _ensure_download_url(file_path: str, download_url: str = "") -> str:
    if download_url:
        return download_url
    if not file_path or not _os.path.exists(file_path):
        return ""
    file_name = _os.path.basename(file_path)
    return f"{_get_public_base_url()}/api/files/download/{file_name}"

def _download_remote_file(file_url: str, allowed_extensions: list = None) -> dict:
    import requests
    if not file_url or not file_url.startswith(("http://", "https://")):
        raise ValueError(f"无效的文件 URL: {file_url}")
    parsed = _urlparse(file_url)
    url_path = _unquote(parsed.path)
    original_name = url_path.split("/")[-1] if "/" in url_path else "downloaded_file"
    ext = ("." + original_name.rsplit(".", 1)[-1].lower()) if "." in original_name else ""
    if allowed_extensions and ext not in allowed_extensions:
        raise ValueError(f"不支持的文件类型: {ext}，允许: {allowed_extensions}")
    short_id = str(_uuid.uuid4())[:8]
    local_name = f"{original_name.rsplit('.', 1)[0]}_{short_id}{ext}" if ext else f"{original_name}_{short_id}"
    save_dir = _get_generated_files_dir()
    save_dir.mkdir(parents=True, exist_ok=True)
    local_path = save_dir / local_name
    resp = requests.get(file_url, timeout=60, stream=True)
    resp.raise_for_status()
    with open(local_path, "wb") as f:
        for chunk in resp.iter_content(chunk_size=8192):
            f.write(chunk)
    download_url = f"{_get_public_base_url()}/api/files/download/{local_name}"
    return {"file_path": str(local_path), "download_url": download_url, "file_name": local_name}
# ──────────────────────────────────────────────────────────────────────────────


# ── 兼容层：SkillResult / SkillStatus（老架构接口，保持向后兼容）──
class _SkillStatus:
    SUCCESS = "success"
    ERROR = "error"
    PARTIAL = "partial"

class _SkillResult(dict):
    """轻量兼容类：SkillResult(status=..., data=..., error=...) 直接作为 dict 使用"""
    def __init__(self, status=None, data=None, error=None, **kwargs):
        d = {}
        if status is not None:
            d["status"] = status
        if data is not None:
            if isinstance(data, dict):
                d.update(data)
            else:
                d["data"] = data
        if error is not None:
            d["error"] = error
        d.update(kwargs)
        super().__init__(d)

SkillResult = _SkillResult
SkillStatus = _SkillStatus
# ────────────────────────────────────────────────────────────────────────────

class PptxPreviewSkill:
    """PowerPoint 演示文稿生成与预览技能"""

    @property
    def name(self) -> str:
        return "pptx_preview"

    @property
    def description(self) -> str:
        return "生成或预览 PowerPoint 演示文稿 (.pptx)。支持从数据生成演示文稿，或预览现有文件，提供下载链接和文本预览。"

    @property
    def category(self) -> str:
        return "office_preview"

    @property
    def input_schema(self) -> Dict[str, Any]:
        return {
            "file_path": {"type": "string", "required": False, "description": "pptx 文件路径 (预览模式)"},
            "data": {
                "type": "object", 
                "required": False, 
                "description": "用于生成演示文稿的数据 (生成模式)。格式: {'title': str, 'subtitle': str, 'slides': [{'title': str, 'content': str|list}]}"
            },
            "file_url": {"type": "string", "required": False, "description": "pptx 文件 URL"},
            "extract_text_only": {"type": "boolean", "default": False, "description": "是否仅提取文本"}
        }

    @property
    def output_schema(self) -> Dict[str, Any]:
        return {
            "file_name": "string",
            "file_size": "integer",
            "file_data_base64": "base64 encoded file content",
            "text_content": "extracted text content",
            "slide_count": "integer",
            "slides_text": "list of slide text content",
            "download_url": "string (optional)"
        }

    async def _extract_text_from_pptx(self, file_path: str) -> tuple[list[str], int]:
        """从 pptx 文件提取文本内容"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_text = []
            slide_count = len(prs.slides)

            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    slides_text.append("\n".join(slide_text))

            return slides_text, slide_count

        except Exception as e:
            print(f"[PptxPreview] Text extraction failed: {e}")
            return [], 0

    async def _read_file_as_base64(self, file_path: str) -> tuple[str, int]:
        """读取文件并转换为 base64"""
        try:
            with open(file_path, 'rb') as f:
                file_data = f.read()
                base64_data = base64.b64encode(file_data).decode('utf-8')
                return base64_data, os.path.getsize(file_path)
        except Exception as e:
            print(f"[PptxPreview] Failed to read file: {e}")
            return "", 0

    async def execute(self, context: dict) -> dict:
        """执行 pptx 生成或预览"""
        try:
            file_path = context.get("file_path", "")
            file_url = context.get("file_url", "")
            data_input = context.get("data", None)
            extract_text_only = context.get("extract_text_only", False)
            download_url = ""

            # 模式 1: 生成模式
            if data_input:
                # 处理 LLM 传入的 JSON 字符串情况
                if isinstance(data_input, str):
                    import json
                    try:
                        data_input = json.loads(data_input)
                        print(f"[PptxPreview] Parsed JSON string to dict: {list(data_input.keys()) if isinstance(data_input, dict) else type(data_input)}")
                    except json.JSONDecodeError as e:
                        return SkillResult(
                            status=SkillStatus.ERROR,
                            error=f"data 参数 JSON 解析失败: {str(e)}"
                        )

                return SkillResult(
                    status=SkillStatus.ERROR,
                    error="请通过 file_url 或 file_path 参数提供文件"
                )
            # 如果没有提供数据输入，根据用户输入和上下文生成动态数据
            elif not file_path and not file_url:
                return SkillResult(
                    status=SkillStatus.ERROR,
                    error="请通过 file_url 或 file_path 参数提供文件"
                )

            # 模式 2: URL 预览模式（下载远程文件到本地）
            elif file_url:
                try:
                    dl_result = _download_remote_file(file_url, allowed_extensions=['.pptx'])
                    file_path = dl_result["file_path"]
                    download_url = dl_result["download_url"]
                    print(f"[PptxPreview] Downloaded remote file: {file_url} -> {file_path}")
                except Exception as e:
                    return SkillResult(
                        status=SkillStatus.ERROR,
                        error=f"远程文件下载失败: {str(e)}"
                    )
            elif not file_path or not os.path.exists(file_path):
                return SkillResult(
                    status=SkillStatus.ERROR,
                    error=f"未提供数据且文件不存在: {file_path}"
                )

            # 验证文件类型
            if not file_path.lower().endswith('.pptx'):
                return SkillResult(
                    status=SkillStatus.ERROR,
                    error="仅支持 .pptx 文件格式"
                )

            file_name = os.path.basename(file_path)

            # 提取文本内容
            slides_text, slide_count = await self._extract_text_from_pptx(file_path)
            text_content = "\n\n---\n\n".join([f"幻灯片 {i+1}\n{text}" for i, text in enumerate(slides_text)])

            # 读取文件为 base64
            file_data_base64, file_size = await self._read_file_as_base64(file_path)

            # 确保有可访问的下载 URL
            if not download_url:
                download_url = _ensure_download_url(file_path, download_url)

            # 构建返回数据
            data = {
                "file_name": file_name,
                "file_size": file_size,
                "file_data_base64": file_data_base64,
                "text_content": text_content,
                "slide_count": slide_count,
                "slides_text": slides_text,
                "slides_data": data_input.get("slides", []) if data_input else [],
                "download_url": download_url,
                "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            }

            return SkillResult(
                status=SkillStatus.SUCCESS,
                data=data,
                message=f"成功处理 PowerPoint 演示文稿: {file_name}"
            )

        except Exception as e:
            import traceback
            traceback.print_exc()
            return SkillResult(
                status=SkillStatus.ERROR,
                error=f"演示文稿处理失败: {str(e)}"
            )

    def _generate_dynamic_pptx_data(self, user_input: str = "", context_data: Dict = None) -> Dict[str, Any]:
        """根据用户输入和上下文动态生成PPTX数据"""

        # 分析用户输入，确定PPT主题
        if not user_input:
            user_input = "生成演示文稿"

        # 从上下文或用户输入中提取数据
        title = self._extract_title_from_input(user_input)

        # 检查用户是否明确要求图表
        wants_charts = any(keyword in user_input.lower() for keyword in
                          ['图表', '图片', 'chart', 'graph', 'visualization', '数据可视化'])

        slides = []

        if wants_charts:
            # 用户要求图表，使用包含图表的默认数据
            print(f"[PptxPreview] User requested charts, using chart data")
            default_data = self._get_default_pptx_data()
            default_data["title"] = title
            default_data["subtitle"] = f"{title} - 深度分析报告"
            return default_data
        elif context_data:
            # 如果有上下文数据，基于数据生成内容
            slides.extend(self._generate_data_based_slides(context_data))
        else:
            # 基于用户输入生成通用内容（也包含图表示例）
            slides.extend(self._generate_generic_slides(user_input))

        return {
            "title": title,
            "subtitle": f"{title} - 深度分析报告",
            "theme": "business",
            "slides": slides
        }

    def _extract_title_from_input(self, user_input: str) -> str:
        """从用户输入中提取标题"""
        # 简单的标题提取逻辑
        keywords = ["生成", "制作", "创建", "PPT", "演示文稿", "分析"]
        cleaned_input = user_input

        for keyword in keywords:
            cleaned_input = cleaned_input.replace(keyword, "")

        cleaned_input = cleaned_input.strip("@pptx_preview ")
        cleaned_input = cleaned_input.strip()

        if not cleaned_input:
            return "通用演示文稿"

        return cleaned_input[:20]  # 限制标题长度

    def _generate_data_based_slides(self, context_data: Dict) -> List[Dict]:
        """基于数据生成幻灯片"""
        slides = []

        # 数据概览幻灯片
        if context_data.get("analysis_data"):
            slides.append({
                "type": "content",
                "title": "数据分析概览",
                "content": [
                    f"数据来源: {context_data.get('source', '未知来源')}",
                    f"数据量: {context_data.get('data_count', 'N/A')}条",
                    f"分析维度: {len(context_data.get('dimensions', []))}个",
                    f"更新时间: {context_data.get('update_time', 'N/A')}"
                ]
            })

        # 动态图表幻灯片
        if context_data.get("chart_data"):
            chart_data = context_data["chart_data"]
            slides.append({
                "type": "chart",
                "title": chart_data.get("title", "数据可视化分析"),
                "content": chart_data.get("description", "数据图表展示"),
                "chart_data": chart_data
            })

        # 图片幻灯片
        if context_data.get("images"):
            for img_data in context_data["images"]:
                slides.append({
                    "type": "image",
                    "title": img_data.get("title", "数据图表"),
                    "image_data": img_data.get("data", ""),
                    "caption": img_data.get("description", "")
                })

        return slides

    def _generate_generic_slides(self, user_input: str) -> List[Dict]:
        """基于用户输入生成通用幻灯片"""
        slides = []
        title = self._extract_title_from_input(user_input)

        # 标题页
        slides.append({
            "type": "title",
            "title": title,
            "content": "演示文稿"
        })

        # 内容介绍页
        slides.append({
            "type": "content",
            "title": "内容概述",
            "content": [
                "本演示文稿基于用户需求自动生成",
                f"主题: {title}",
                "内容结构清晰，重点突出",
                "可根据需要进一步定制化"
            ]
        })

        # 图表示例页（生成真实的图表）
        slides.append({
            "type": "chart",
            "title": f"{title}数据可视化",
            "content": "示例数据图表展示",
            "chart": {
                "type": "bar",
                "title": "示例柱状图",
                "categories": ["类别A", "类别B", "类别C", "类别D", "类别E"],
                "values": [25, 45, 30, 60, 40],  # 示例数据
                "colors": ["FF6B6B", "4ECDC4", "45B7D1", "96CEB4", "FFEAA7"]
            }
        })

        # 饼图示例
        slides.append({
            "type": "chart",
            "title": f"{title}分布分析",
            "content": "数据分布饼图展示",
            "chart": {
                "type": "pie",
                "title": "示例饼图",
                "categories": ["技术A", "技术B", "技术C", "技术D"],
                "values": [35, 25, 25, 15],
                "colors": ["FF6B6B", "4ECDC4", "45B7D1", "96CEB4"]
            }
        })

        # 结论页
        slides.append({
            "type": "content",
            "title": "总结",
            "content": [
                "感谢观看本次演示",
                "支持真实的图表和数据可视化",
                "可根据实际数据动态生成图表",
                "欢迎提供更多详细信息进行定制"
            ]
        })

        return slides

    def _generate_minimal_pptx_data(self, user_input: str = "") -> Dict[str, Any]:
        """生成最简单的PPTX数据作为降级方案"""
        title = self._extract_title_from_input(user_input) or "简单演示文稿"

        return {
            "title": title,
            "subtitle": f"{title} - 基础版本",
            "theme": "business",
            "slides": [
                {
                    "type": "title",
                    "title": title,
                    "content": "基础演示文稿"
                },
                {
                    "type": "content",
                    "title": "内容说明",
                    "content": [
                        "这是一个基础版本的演示文稿",
                        "由于数据处理遇到问题，使用简化的内容",
                        "如果您需要具体的内容，请提供更多信息",
                        "您可以尝试重新生成或提供具体的数据"
                    ]
                }
            ]
        }

    def _extract_data_from_context(self, context: dict) -> Dict:
        """从上下文中提取可能的数据"""
        data_sources = []

        # 尝试从shared_data中获取数据
        if hasattr(context, 'shared_data') and context.shared_data:
            data_sources.append(context.shared_data)

        # 尝试从session_data中获取
        if hasattr(context, 'session_data') and context.session_data:
            data_sources.append(context.session_data)

        # 合并所有数据源
        merged_data = {}
        for source in data_sources:
            if isinstance(source, dict):
                merged_data.update(source)

        return merged_data

    def _get_default_pptx_data(self) -> Dict[str, Any]:
        """获取默认的PPTX数据结构，包含图表"""
        return {
            "title": "数据可视化演示",
            "subtitle": "数据可视化演示 - 深度分析报告",
            "theme": "business",
            "slides": [
                {
                    "type": "title",
                    "title": "数据可视化演示",
                    "content": "演示文稿"
                },
                {
                    "type": "content",
                    "title": "概览",
                    "content": [
                        "这是一个包含多种图表类型的演示文稿",
                        "支持柱状图、饼图、折线图、环形图",
                        "新增思维导图和时间轴图表"
                    ]
                },
                {
                    "type": "chart",
                    "title": "柱状图示例",
                    "content": "数据展示",
                    "chart": {
                        "type": "bar",
                        "title": "示例柱状图",
                        "categories": ["类别A", "类别B", "类别C", "类别D", "类别E"],
                        "values": [25, 45, 30, 60, 40],
                        "colors": ["FF6B6B", "4ECDC4", "45B7D1", "96CEB4", "FFEAA7"]
                    }
                },
                {
                    "type": "chart",
                    "title": "折线图示例 - 时间戳数据",
                    "content": "趋势分析",
                    "chart": {
                        "type": "line",
                        "title": "市场增长趋势",
                        "x": [1609459200, 1640995200, 1672531200, 1704067200],  # 2021-01-01, 2022-01-01, 2023-01-01, 2024-01-01
                        "y": [100, 150, 220, 350],
                        "colors": ["4ECDC4"],
                        "x_label": "时间",
                        "y_label": "市场规模（亿元）"
                    }
                },
                {
                    "type": "chart",
                    "title": "饼图示例",
                    "content": "数据分布",
                    "chart": {
                        "type": "pie",
                        "title": "示例饼图",
                        "categories": ["技术A", "技术B", "技术C", "技术D"],
                        "values": [35, 25, 25, 15],
                        "colors": ["FF6B6B", "4ECDC4", "45B7D1", "96CEB4"]
                    }
                },
                {
                    "type": "chart",
                    "title": "环形图示例",
                    "content": "占比分析",
                    "chart": {
                        "type": "doughnut",
                        "title": "市场份额分析",
                        "categories": ["产品A", "产品B", "产品C", "产品D"],
                        "values": [40, 30, 20, 10],
                        "colors": ["FF6B6B", "4ECDC4", "45B7D1", "96CEB4"]
                    }
                },
                {
                    "type": "chart",
                    "title": "思维导图示例",
                    "content": "架构分析",
                    "chart": {
                        "type": "mindmap",
                        "title": "技术架构思维导图",
                        "center": "核心技术",
                        "center_size": 120,
                        "nodes": ["前端技术", "后端架构", "数据存储", "AI算法"],
                        "values": [80, 100, 60, 90],
                        "colors": ["FF6B6B", "4ECDC4", "45B7D1", "96CEB4"]
                    }
                },
                {
                    "type": "chart",
                    "title": "时间轴示例",
                    "content": "发展历程",
                    "chart": {
                        "type": "timeline",
                        "title": "产品发展时间轴",
                        "years": ["2020", "2021", "2022", "2023"],
                        "values": [50, 120, 280, 450],
                        "colors": ["FF6B6B", "4ECDC4", "45B7D1", "96CEB4"],
                        "x_label": "用户数量（万）"
                    }
                },
                {
                    "type": "content",
                    "title": "总结",
                    "content": [
                        "感谢观看",
                        "图表已成功生成",
                        "支持的数据可视化"
                    ]
                }
            ]
        }


def _main():
    """直接执行入口: python3 script.py --param1 value1
    也支持 JSON stdin: echo '{"param1": "v1"}' | python3 script.py
    """
    import argparse
    import asyncio
    import json
    import sys

    params = {}
    if not sys.stdin.isatty():
        try:
            raw = sys.stdin.read().strip()
            if raw:
                params = json.loads(raw)
        except Exception:
            pass

    parser = argparse.ArgumentParser(description="Run PptxPreviewSkill directly")
    parser.add_argument("--file-path", type=str, dest="file_path")
    parser.add_argument("--data", type=str, dest="data")
    args = parser.parse_args()
    for k, v in vars(args).items():
        if v is not None:
            params[k] = v

    async def run():
        skill = PptxPreviewSkill()
        result = await skill.execute(params)
        out = result if isinstance(result, dict) else {"data": str(result)}
        print(json.dumps(out, ensure_ascii=False, default=str, indent=2))

    asyncio.run(run())


if __name__ == "__main__":
    _main()
